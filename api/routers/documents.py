"""Document-centric API endpoints."""

from __future__ import annotations

import logging
import os
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, Dict, List, Optional, Tuple

import pdfplumber
from botocore.exceptions import ClientError
from docx import Document as DocxDocument
from fastapi import APIRouter, Depends, File, Form, HTTPException, Request, UploadFile
from pptx import Presentation
from pydantic import BaseModel, Field, field_validator

from services.document_embedding_service import DocumentEmbeddingService
from services.document_extractor import DocumentExtractor
from services.model_selector import RAGPipeline

# Ensure GPU-related environment variables align with the rest of the API.
os.environ.setdefault("CUDA_VISIBLE_DEVICES", "0")
os.environ.setdefault("OLLAMA_USE_GPU", "1")
os.environ.setdefault("OLLAMA_NUM_PARALLEL", "4")
os.environ.setdefault("OMP_NUM_THREADS", "8")


logger = logging.getLogger(__name__)

router = APIRouter(prefix="/document", tags=["Documents"])


DEFAULT_S3_BUCKET_NAME = "procwisemvp"
SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".txt"}


def get_agent_nick(request: Request):
    agent_nick = getattr(request.app.state, "agent_nick", None)
    if not agent_nick:
        raise HTTPException(status_code=503, detail="AgentNick not available")
    return agent_nick


def get_rag_pipeline(request: Request) -> RAGPipeline:
    pipeline = getattr(request.app.state, "rag_pipeline", None)
    if not pipeline:
        raise HTTPException(
            status_code=503, detail="RAG Pipeline service is not available."
        )
    return pipeline


def _extract_text_from_bytes(filename: str, data: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported file type '{suffix or 'unknown'}'. "
                "Allowed types: PDF, DOCX, PPTX, TXT."
            ),
        )

    if suffix == ".pdf":
        try:
            with pdfplumber.open(BytesIO(data)) as pdf:
                pages = [page.extract_text() for page in pdf.pages]
        except Exception as exc:  # pragma: no cover - depends on document layout
            raise HTTPException(
                status_code=400,
                detail=f"Failed to read PDF document '{filename}': {exc}",
            ) from exc
        return "\n".join(page for page in pages if page)

    if suffix == ".docx":
        try:
            document = DocxDocument(BytesIO(data))
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to read DOCX document '{filename}': {exc}",
            ) from exc
        return "\n".join(
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text and paragraph.text.strip()
        )

    if suffix == ".pptx":
        try:
            presentation = Presentation(BytesIO(data))
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail=f"Failed to read PPTX document '{filename}': {exc}",
            ) from exc
        slide_text: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
        return "\n".join(slide_text)

    try:
        return data.decode("utf-8-sig")
    except UnicodeDecodeError:
        return data.decode("utf-8", errors="ignore")


class S3DocumentExtractionRequest(BaseModel):
    """Input payload for S3-backed document extraction."""

    s3_path: str = Field(
        ...,
        description="S3 folder or object path containing procurement documents",
    )
    metadata: Dict[str, Any] = Field(
        default_factory=dict,
        description="Optional metadata applied to every ingested document",
    )

    @field_validator("s3_path", mode="before")
    @classmethod
    def _strip_s3_path(cls, value: Any) -> str:
        if value is None:
            raise ValueError("s3_path is required")
        if not isinstance(value, str):
            value = str(value)
        path = value.strip()
        if not path:
            raise ValueError("s3_path must not be empty")
        return path


class DocumentEmbeddingResponse(BaseModel):
    """Response payload for embedded document uploads."""

    document_id: str
    collection: str
    chunk_count: int
    metadata: Dict[str, Any]


class DocumentEmbeddingError(BaseModel):
    """Details about documents that failed to embed."""

    filename: str
    reason: str


class DocumentEmbeddingBatchResponse(BaseModel):
    """Aggregate response for multi-document embedding uploads."""

    status: str
    total_documents: int
    total_chunks: int
    processed: List[DocumentEmbeddingResponse]
    failed: List[DocumentEmbeddingError] = Field(default_factory=list)


def _resolve_s3_path(s3_path: str, default_bucket: str) -> Tuple[str, str]:
    """Return bucket and prefix for the provided path."""

    value = s3_path.strip()
    if not value:
        raise ValueError("s3_path must not be empty")

    if value.startswith("s3://"):
        stripped = value[5:]
        bucket_part, _, key_part = stripped.partition("/")
        bucket_name = bucket_part or default_bucket
        prefix = key_part
    else:
        bucket_name = default_bucket
        prefix = value

    prefix = prefix.lstrip("/")
    if not prefix:
        raise ValueError("s3_path must include an object prefix")
    return bucket_name, prefix


def _list_s3_keys(client: Any, bucket: str, prefix: str) -> List[str]:
    """List all object keys for the provided prefix."""

    keys: List[str] = []
    continuation: Optional[str] = None

    while True:
        params: Dict[str, Any] = {"Bucket": bucket, "Prefix": prefix}
        if continuation:
            params["ContinuationToken"] = continuation
        response = client.list_objects_v2(**params)
        contents = (response or {}).get("Contents", [])
        for entry in contents:
            key = entry.get("Key") if isinstance(entry, dict) else None
            if key and not key.endswith("/"):
                keys.append(key)

        if not response or not response.get("IsTruncated"):
            break
        continuation = response.get("NextContinuationToken")
        if not continuation:
            break

    return keys


def _download_s3_object(client: Any, bucket: str, key: str) -> bytes:
    """Download an S3 object and return its bytes."""

    response = client.get_object(Bucket=bucket, Key=key)
    body = response.get("Body") if isinstance(response, dict) else None
    if body is None:
        raise HTTPException(status_code=500, detail=f"S3 object body was empty for {key}")

    try:
        return body.read()
    finally:
        try:
            body.close()
        except Exception:  # pragma: no cover - defensive cleanup
            logger.debug("Failed to close S3 streaming body for %s", key, exc_info=True)


@router.post("/extract-from-s3")
def extract_document_from_s3(
    req: S3DocumentExtractionRequest,
    agent_nick=Depends(get_agent_nick),
):
    """Download a document from S3, extract it, and persist the structured payload."""

    bucket = getattr(agent_nick.settings, "s3_bucket_name", None)
    if not bucket:
        bucket = os.getenv("S3_BUCKET_NAME", DEFAULT_S3_BUCKET_NAME)

    if not bucket:
        raise HTTPException(status_code=503, detail="S3 bucket is not configured")

    try:
        effective_bucket, prefix = _resolve_s3_path(req.s3_path, bucket)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    connection_factory = lambda: agent_nick.get_db_connection()
    chat_options = {}
    if hasattr(agent_nick, "ollama_options") and callable(agent_nick.ollama_options):
        try:
            chat_options = agent_nick.ollama_options()
        except Exception:  # pragma: no cover - fall back to defaults
            logger.debug("Failed to obtain Ollama chat options", exc_info=True)
            chat_options = {}

    extractor = DocumentExtractor(
        connection_factory,
        preferred_models=("qwen3", "phi4"),
        chat_options=chat_options,
    )

    base_metadata = dict(req.metadata or {})

    documents: List[Dict[str, Any]] = []

    try:
        with agent_nick.reserve_s3_connection() as s3_client:
            try:
                object_keys = _list_s3_keys(s3_client, effective_bucket, prefix)
            except ClientError as exc:  # pragma: no cover - network failures in prod
                error = exc.response.get("Error", {}) if hasattr(exc, "response") else {}
                detail = error.get("Message") or f"Unable to list objects for {prefix}"
                raise HTTPException(status_code=500, detail=detail) from exc
            except Exception as exc:  # pragma: no cover - defensive
                logger.exception("Failed to list S3 objects for %s", prefix)
                raise HTTPException(status_code=500, detail="Failed to list S3 objects") from exc

            if not object_keys:
                raise HTTPException(status_code=404, detail="No documents found for provided s3_path")

            for object_key in object_keys:
                if not object_key:
                    continue
                suffix = os.path.splitext(object_key)[1] or ".bin"
                try:
                    file_bytes = _download_s3_object(s3_client, effective_bucket, object_key)
                except ClientError as exc:  # pragma: no cover - network failures in prod
                    error = exc.response.get("Error", {}) if hasattr(exc, "response") else {}
                    code = (error or {}).get("Code")
                    status = 404 if code in {"NoSuchKey", "404"} else 500
                    detail = error.get("Message") or f"Unable to download {object_key}"
                    raise HTTPException(status_code=status, detail=detail) from exc
                except HTTPException:
                    raise
                except Exception as exc:  # pragma: no cover - defensive
                    logger.exception("Failed to download %s from S3", object_key)
                    raise HTTPException(status_code=500, detail="Failed to download document from S3") from exc

                tmp_path = None
                try:
                    with NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = Path(tmp.name)

                    metadata = dict(base_metadata)
                    metadata.setdefault("ingestion_mode", "s3")
                    metadata.setdefault("s3_bucket", effective_bucket)
                    metadata.setdefault("s3_object_key", object_key)

                    result = extractor.extract(
                        tmp_path,
                        metadata=metadata,
                        source_label=object_key,
                    )
                except ValueError as exc:
                    raise HTTPException(status_code=400, detail=str(exc)) from exc
                except HTTPException:
                    raise
                except Exception as exc:  # pragma: no cover - defensive
                    logger.exception("Document extraction failed for %s", object_key)
                    raise HTTPException(status_code=500, detail="Document extraction failed") from exc
                finally:
                    if tmp_path:
                        try:
                            tmp_path.unlink()
                        except FileNotFoundError:
                            pass
                        except Exception:  # pragma: no cover - defensive cleanup
                            logger.debug("Failed to remove temporary file %s", tmp_path, exc_info=True)

                documents.append(result.to_json())
    except HTTPException:
        raise
    except Exception as exc:  # pragma: no cover - defensive cleanup
        logger.exception("Unexpected failure during S3 extraction")
        raise HTTPException(status_code=500, detail="Document extraction failed") from exc

    response_payload: Dict[str, Any] = {"status": "completed", "documents": documents}
    if documents:
        response_payload["document"] = documents[0]
    return response_payload


@router.post("/embed-document", response_model=DocumentEmbeddingBatchResponse)
async def embed_documents(
    request: Request,
    files: List[UploadFile] = File(...),
    user_id: Optional[str] = Form(None),
    pipeline: RAGPipeline = Depends(get_rag_pipeline),
    agent_nick=Depends(get_agent_nick),
):
    """Upload, extract, embed, and register user-provided documents for the RAG pipeline."""

    if not files:
        raise HTTPException(status_code=400, detail="At least one document must be provided")

    def _clean(value: Optional[str]) -> Optional[str]:
        if value is None:
            return None
        if not isinstance(value, str):
            value = str(value)
        value = value.strip()
        return value or None

    header_user = _clean(request.headers.get("x-user-id"))
    header_session = _clean(request.headers.get("x-session-id"))
    query_user = _clean(request.query_params.get("user_id"))

    resolved_user: Optional[str] = None
    for candidate in (_clean(user_id), header_user, query_user):
        if candidate:
            resolved_user = candidate
            break

    rag_service = pipeline.rag
    collection_name = getattr(rag_service, "uploaded_collection", "uploaded_documents")

    try:
        rag_service.ensure_collection(collection_name)
    except HTTPException:
        raise
    except Exception as exc:  # pragma: no cover - defensive
        logger.exception("Failed to initialise Qdrant collection for uploads")
        raise HTTPException(
            status_code=500,
            detail="Failed to initialise vector collection for uploaded documents.",
        ) from exc

    embedding_service = DocumentEmbeddingService(agent_nick, collection_name=collection_name)

    processed: List[DocumentEmbeddingResponse] = []
    failures: List[DocumentEmbeddingError] = []
    total_chunks = 0

    for upload in files:
        filename = Path(upload.filename or "uploaded_document").name
        suffix = Path(filename).suffix.lower()
        if suffix not in SUPPORTED_SUFFIXES:
            failures.append(
                DocumentEmbeddingError(
                    filename=filename,
                    reason=(
                        f"Unsupported file type '{suffix or 'unknown'}'. "
                        "Allowed types: PDF, DOCX, PPTX, TXT."
                    ),
                )
            )
            continue

        try:
            data = await upload.read()
        except Exception as exc:
            failures.append(
                DocumentEmbeddingError(
                    filename=filename,
                    reason=f"Unable to read file: {exc}",
                )
            )
            continue

        if not data:
            failures.append(
                DocumentEmbeddingError(
                    filename=filename,
                    reason="Uploaded file contained no data.",
                )
            )
            continue

        metadata: Dict[str, Any] = {
            "mime_type": upload.content_type or "",
            "ingestion_source": "document_embed_endpoint",
        }
        if resolved_user:
            metadata["uploaded_by"] = resolved_user

        try:
            embedded = embedding_service.embed_document(
                filename=filename,
                file_bytes=data,
                metadata=metadata,
            )
        except ValueError as exc:
            failures.append(DocumentEmbeddingError(filename=filename, reason=str(exc)))
            continue
        except HTTPException:
            raise
        except Exception as exc:  # pragma: no cover - defensive
            logger.exception("Embedding pipeline failed for %s", filename)
            failures.append(
                DocumentEmbeddingError(
                    filename=filename,
                    reason="Unexpected error while embedding document.",
                )
            )
            continue

        processed.append(
            DocumentEmbeddingResponse(
                document_id=embedded.document_id,
                collection=embedded.collection,
                chunk_count=embedded.chunk_count,
                metadata=embedded.metadata,
            )
        )
        total_chunks += embedded.chunk_count

    if not processed:
        detail = failures[0].reason if failures else "No documents were processed."
        raise HTTPException(status_code=400, detail=detail)

    uploaded_document_ids = [doc.document_id for doc in processed]
    upload_metadata = {
        "filenames": [
            doc.metadata.get("filename") or doc.metadata.get("doc_name") for doc in processed
        ],
        "total_chunks": total_chunks,
    }
    if resolved_user:
        upload_metadata["uploaded_by"] = resolved_user

    try:
        pipeline.activate_uploaded_context(
            uploaded_document_ids,
            metadata=upload_metadata,
            session_id=header_session or resolved_user,
        )
    except AttributeError:
        logger.debug(
            "RAG pipeline does not expose uploaded context activation",
            exc_info=True,
        )

    return DocumentEmbeddingBatchResponse(
        status="success",
        total_documents=len(processed),
        total_chunks=total_chunks,
        processed=processed,
        failed=failures,
    )


__all__ = ["router"]
